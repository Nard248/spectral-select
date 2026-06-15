"""Build the PowerPoint deck for Narek Meloyan's Master's Thesis defense.

This version (v2, 2026-05-14) mirrors the narrative arc requested by Narek:
broad vision-biology framing -> general HSI -> biomedicine -> fluorescence
physics -> single-excitation HSI -> its limitations -> ME-HSI as the solution
-> ME-HSI's own challenges -> our framework -> method -> datasets -> results
-> validation -> conclusion.

Output: MasterThesis_Narek_Meloyan_Defense.pptx (16:9, ~30 slides, full
speaker notes).

Run from project root:
    .venv/bin/python MasterThesis_Narek_Meloyan/build_defense_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent / "MasterThesis_Narek_Meloyan_Defense.pptx"
IMG3P = Path(__file__).resolve().parent / "images_third_party"

# Theme colors
NAVY = RGBColor(0x1F, 0x4E, 0x79)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0x88, 0x88, 0x88)
ACCENT_RED = RGBColor(0xC8, 0x10, 0x2E)
ACCENT_GREEN = RGBColor(0x2E, 0x8B, 0x57)
ACCENT_BLUE = RGBColor(0x1F, 0x77, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAINT = RGBColor(0xE8, 0xEE, 0xF6)
PALE = RGBColor(0xF4, 0xF6, 0xFA)


# ===========================================================================
# Helpers
# ===========================================================================
def add_text_box(slide, left, top, width, height, text, *,
                 font_size=18, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, italic=False, name=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if name:
        r.font.name = name
    return tb


def add_bullet_list(slide, left, top, width, height, bullets, *,
                    font_size=18, color=DARK_GREY, level_size=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        indent = 0
        text = line
        if line.startswith("    - "):
            indent = 2
            text = line[6:]
        elif line.startswith("  - "):
            indent = 1
            text = line[4:]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        r = p.add_run()
        marker = "• " if indent == 0 else "– "
        r.text = marker + text
        r.font.size = Pt(font_size if indent == 0 else (level_size or font_size - 3))
        r.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(13.333), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6),
                 title, font_size=26, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
    return bar


def add_footer(slide, slide_number, total):
    add_text_box(slide, Inches(0.4), Inches(7.10), Inches(7), Inches(0.3),
                 "Master's Thesis Defense — Narek Meloyan — 2026",
                 font_size=9, color=LIGHT_GREY, italic=True)
    add_text_box(slide, Inches(12.2), Inches(7.10), Inches(0.8), Inches(0.3),
                 f"{slide_number} / {total}",
                 font_size=9, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left, top, width or Inches(4), height or Inches(3))
        ph.fill.solid()
        ph.fill.fore_color.rgb = FAINT
        ph.line.color.rgb = LIGHT_GREY
        add_text_box(slide, left, top + Inches(1), width or Inches(4), Inches(0.5),
                     f"[Figure: {Path(path).name}]",
                     font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        return None
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def add_caption(slide, left, top, width, text):
    add_text_box(slide, left, top, width, Inches(0.4), text,
                 font_size=11, italic=True, color=LIGHT_GREY,
                 align=PP_ALIGN.CENTER)


def add_chip(slide, left, top, width, height, text, *,
             fill=FAINT, text_color=NAVY, font_size=14, bold=True):
    """A colored rounded rectangle with centered text — used for concept tags."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = NAVY
    s.line.width = Pt(0.75)
    add_text_box(slide, left, top, width, height, text,
                 font_size=font_size, bold=bold, color=text_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def set_notes(slide, notes_text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = notes_text


def add_image_credit(slide, left, top, width, text):
    """Tiny grey attribution line for third-party images."""
    add_text_box(slide, left, top, width, Inches(0.3),
                 f"Image: {text}",
                 font_size=8, italic=True, color=LIGHT_GREY,
                 align=PP_ALIGN.LEFT)


# ===========================================================================
# Build the deck
# ===========================================================================
def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    SLIDES_PLANNED = 31
    def new_slide():
        return prs.slides.add_slide(blank)

    # ---------- 1. TITLE ----------
    s = new_slide()
    # Decorative HSI cube image on the left
    add_image(s, IMG3P / "hsi_cube.png", Inches(0.5), Inches(1.6),
              width=Inches(3.2))
    # Decorative fluorescent minerals image on the right
    add_image(s, IMG3P / "fluorescent_minerals.jpg", Inches(9.7), Inches(1.6),
              width=Inches(3.2))
    add_text_box(s, Inches(4.0), Inches(1.6), Inches(5.3), Inches(1.2),
                 "Deep Learning for\nDimensionality Reduction in",
                 font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(4.0), Inches(2.9), Inches(5.3), Inches(1.0),
                 "Multi-Excitation\nHyperspectral Imaging",
                 font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                 "Master's Thesis Defense",
                 font_size=22, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
                 "Narek Meloyan",
                 font_size=24, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
                 "American University of Armenia · 2026",
                 font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                 "Supervisor: Prof. Narine Sarvazyan",
                 font_size=13, italic=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_image_credit(s, Inches(0.5), Inches(6.9), Inches(12.3),
                     "Wikimedia Commons (CC BY-SA) — see credits slide")
    set_notes(s,
        "Good afternoon. My name is Narek Meloyan, and today I'll present my "
        "Master's thesis on deep learning for dimensionality reduction in "
        "multi-excitation hyperspectral imaging. This work was supervised by "
        "Professor Narine Sarvazyan. Over the next ~30 minutes I'll start with a "
        "broad introduction to hyperspectral imaging, walk through how "
        "fluorescence imaging works, explain why we need the 4D extension I "
        "studied, present the deep-learning framework I built, and show "
        "experimental results on two biological fluorescence datasets.")

    # ---------- 2. OUTLINE ----------
    s = new_slide()
    add_title_bar(s, "Outline")
    add_bullet_list(s, Inches(1.5), Inches(1.4), Inches(10), Inches(5.5), [
        "1. From human vision to hyperspectral imaging",
        "2. HSI in the biomedical setting",
        "3. Fluorescence: what it is, how we capture it",
        "4. Why single-excitation fluorescence HSI is not enough",
        "5. ME-HSI: the 4D extension, and its own challenges",
        "6. Existing band-selection methods and why they fall short",
        "7. The proposed framework: 3D CAE + perturbation + MMR",
        "8. Two datasets: Lichens and Collagen Sponges",
        "9. Results & validation",
        "10. Discussion, limitations, conclusion",
    ], font_size=20)
    add_footer(s, 2, SLIDES_PLANNED)
    set_notes(s,
        "Here's the plan. The first six slides build the context — what HSI is "
        "and why we care about ME-HSI in the biomedical context. The middle "
        "third explains the method I developed. The final third presents "
        "experimental results on two datasets and the validation. I'll keep the "
        "intro brisk so we spend most of the time on the contribution and "
        "results.")

    # ---------- 3. VISION IN NATURE ----------
    s = new_slide()
    add_title_bar(s, "How we see — and how we don't")
    # 4 image+label tiles across the top half of the slide
    tiles = [
        ("honeybee.jpg",          "Honeybee",
         "3 cones (UV, blue, green)\nsees UV patterns",                    ACCENT_GREEN, Inches(0.4)),
        ("monarch_butterfly.jpg", "Butterfly",
         "5 – 9 photoreceptors\nbroad spectral range",                     ACCENT_RED,   Inches(3.65)),
        ("mantis_full.jpg",       "Mantis shrimp",
         "12 – 16 photoreceptors\n+ polarisation",                         ACCENT_BLUE,  Inches(6.90)),
        ("hsi_cube_alt.jpg",      "Hyperspectral camera",
         "100 – 300+ channels\ncontinuous spectrum",                       NAVY,         Inches(10.15)),
    ]
    img_w = Inches(3.05); img_h = Inches(2.3)
    for fname, label, body, col, left in tiles:
        add_image(s, IMG3P / fname, left, Inches(1.15), width=img_w, height=img_h)
        # Coloured rule under each image
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.50),
                                  img_w, Inches(0.05))
        rule.fill.solid(); rule.fill.fore_color.rgb = col; rule.line.fill.background()
        add_text_box(s, left, Inches(3.60), img_w, Inches(0.4), label,
                     font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text_box(s, left, Inches(4.05), img_w, Inches(1.1), body,
                     font_size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
    # Bottom: spectrum comparison strip
    add_image(s, IMG3P / "em_spectrum.jpg", Inches(2.0), Inches(5.2),
              width=Inches(9.3))
    add_text_box(s, Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.4),
                 "More channels → more discrimination. A hyperspectral camera is to "
                 "colour vision what colour vision is to greyscale.",
                 font_size=12, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_image_credit(s, Inches(0.4), Inches(7.05), Inches(12.5),
                     "Photos: Wikimedia Commons CC BY-SA — see credits slide")
    add_footer(s, 3, SLIDES_PLANNED)
    set_notes(s,
        "To anchor what hyperspectral imaging is, it helps to start with how "
        "different organisms see. Humans see in three channels — red, green, "
        "blue — and that's why a digital camera also has three. Honeybees see in "
        "three channels too, but shifted into the UV; flowers that look uniform "
        "to us have intricate UV nectar guides bees can see. Butterflies have "
        "around 9 photoreceptor types, and the mantis shrimp famously has "
        "12-to-16 — plus it can detect light polarisation. Each step up means "
        "more discrimination — you can tell apart materials that looked "
        "identical at the lower channel count. A hyperspectral camera takes "
        "this all the way: 100 to 300 contiguous spectral channels, sampling "
        "the full visible-to-near-infrared spectrum continuously rather than "
        "binning it into 3 broad bands. That's the basic 'why' of HSI.")

    # ---------- 4. WHAT HSI BUYS YOU ----------
    s = new_slide()
    add_title_bar(s, "What hyperspectral imaging buys you")
    # Left column: applications list
    add_text_box(s, Inches(0.4), Inches(1.1), Inches(5.3), Inches(0.5),
                 "Applications across science and industry",
                 font_size=16, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.4), Inches(1.6), Inches(5.3), Inches(5.5), [
        "Remote sensing: mineralogy, vegetation, water quality",
        "Agriculture: ripeness, disease, soil moisture",
        "Food industry: contamination, freshness, fraud",
        "Art conservation: pigment ID, hidden layers",
        "Pharmaceutical: powder homogeneity, tablet QA",
        "Defence: camouflage detection",
        "Forensics: trace evidence, fingerprints",
        "Biomedicine: tissue classification, surgical guidance",
    ], font_size=13)
    # Middle: example — false-colour Landsat image
    add_image(s, IMG3P / "landsat_ethiopia.jpg", Inches(5.9), Inches(1.1),
              width=Inches(3.6))
    add_caption(s, Inches(5.9), Inches(4.7), Inches(3.6),
                "Landsat-5 false-colour (Ethiopia)\nVegetation = red, water = blue")
    # Right: why-it-works box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(9.7), Inches(1.1), Inches(3.4), Inches(5.6))
    box.fill.solid(); box.fill.fore_color.rgb = PALE
    box.line.color.rgb = NAVY; box.line.width = Pt(1.2)
    add_text_box(s, Inches(9.85), Inches(1.25), Inches(3.15), Inches(0.5),
                 "Why all of these?", font_size=15, bold=True, color=NAVY)
    add_text_box(s, Inches(9.85), Inches(1.8), Inches(3.15), Inches(4.7),
                 "Materials have characteristic spectral fingerprints.\n\n"
                 "Two surfaces that look identical to the eye can be told apart "
                 "by their full spectrum.\n\n"
                 "HSI gives a per-pixel spectrum — a 'mini-lab' at every point.",
                 font_size=12, color=DARK_GREY)
    add_image_credit(s, Inches(0.4), Inches(7.05), Inches(12.5),
                     "Landsat-5 false-colour, NASA / USGS, public domain (Wikimedia)")
    add_footer(s, 4, SLIDES_PLANNED)
    set_notes(s,
        "Once you have hundreds of channels per pixel, HSI shows up everywhere. "
        "Geologists use it to identify minerals from satellite data. Farmers use "
        "it to spot crop disease before it's visible. The food industry uses it "
        "to detect contamination on processing lines. Art conservators use it "
        "to identify pigments and reveal underdrawings beneath painted "
        "surfaces. The defence community uses it for camouflage and target "
        "identification. And — what this thesis cares about — biomedicine. The "
        "common thread is that materials carry characteristic spectral "
        "fingerprints. Two surfaces that look identical to the human eye can be "
        "told apart by their full spectrum. HSI essentially gives you a "
        "miniature laboratory measurement at every pixel.")

    # ---------- 5. WHY BIOMEDICINE? ----------
    s = new_slide()
    add_title_bar(s, "Why biomedicine cares about HSI")
    add_bullet_list(s, Inches(0.4), Inches(1.2), Inches(8.4), Inches(5.5), [
        "Conventional imaging (RGB, MRI, CT) tells you WHERE things are — not WHAT they are made of",
        "Biological tissues differ chemically: healthy vs. diseased, malignant vs. benign, well-perfused vs. ischemic",
        "These chemical differences are invisible in RGB but visible in the spectrum",
        "Clinical use cases:",
        "  - Tumour-margin assessment during surgery",
        "  - Real-time tissue classification in endoscopy",
        "  - Wound assessment, burn depth estimation",
        "  - Skin-cancer screening, retinal disease",
        "  - Histopathology label-free analysis",
        "Goal: label-free, non-invasive characterisation of tissue in vivo",
    ], font_size=14)
    # Right column: cell-imaging photograph
    add_image(s, IMG3P / "fluorescence_micro.jpg", Inches(9.05), Inches(1.3),
              width=Inches(4.0))
    add_caption(s, Inches(9.05), Inches(4.4), Inches(4.0),
                "Fluorescence microscopy:\nlabel-free contrast at the cellular scale")
    add_image_credit(s, Inches(0.4), Inches(7.05), Inches(12.5),
                     "GFP super-resolution micrograph: C. Cremer, Wikimedia CC BY-SA")
    add_footer(s, 5, SLIDES_PLANNED)
    set_notes(s,
        "Biomedicine cares about HSI because conventional medical imaging — RGB, "
        "MRI, CT — tells you where structures are but not what they're made of. "
        "Biological tissues differ chemically: healthy vs diseased, malignant "
        "vs benign, well-perfused vs ischemic. Those chemical differences are "
        "often invisible in RGB but visible in the full spectrum. Use cases "
        "include tumour margin assessment during surgery, where the surgeon "
        "needs to know in real time whether to cut here or there; endoscopic "
        "tissue classification; burn-depth estimation; skin cancer screening; "
        "and label-free histopathology. The overarching goal is non-invasive, "
        "label-free characterisation of tissue, in vivo.")

    # ---------- 6. FLUORESCENCE PHYSICS ----------
    s = new_slide()
    add_title_bar(s, "Fluorescence: how molecules emit light")
    # Left: Jablonski diagram
    add_image(s, IMG3P / "jablonski.png", Inches(0.4), Inches(1.1),
              height=Inches(4.6))
    add_caption(s, Inches(0.3), Inches(5.8), Inches(3.4),
                "Jablonski diagram of fluorescence")
    # Middle: 5-step process
    add_text_box(s, Inches(3.9), Inches(1.1), Inches(4.5), Inches(0.5),
                 "The physical process", font_size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(3.9), Inches(1.6), Inches(4.5), Inches(4.0), [
        "1. High-energy photon absorbed",
        "2. Electron jumps to excited state",
        "3. Some energy lost as heat",
        "4. Electron drops back, emits new photon",
        "5. λ_emission > λ_excitation (Stokes shift)",
    ], font_size=12)
    # Right: real-world demo via fluorescent minerals
    add_image(s, IMG3P / "fluorescent_minerals.jpg", Inches(8.55), Inches(1.1),
              width=Inches(4.6))
    add_caption(s, Inches(8.55), Inches(4.05), Inches(4.6),
                "Real minerals under UV — each fluorophore\nemits a different colour")
    # Bottom: biological fluorophores band
    add_text_box(s, Inches(3.9), Inches(5.5), Inches(4.5), Inches(0.4),
                 "Biological fluorophores:", font_size=12, bold=True, color=NAVY)
    add_bullet_list(s, Inches(3.9), Inches(5.85), Inches(4.5), Inches(1.3), [
        "NADH 340/460 · Collagen 330/400",
        "Elastin 325/410 · Tryptophan 280/350",
        "Flavins 450/525 · Chlorophyll 430/680",
    ], font_size=10)
    add_image_credit(s, Inches(0.4), Inches(7.05), Inches(12.5),
                     "Jablonski diagram + fluorescent minerals (H. Grobe): Wikimedia CC BY-SA")
    add_footer(s, 6, SLIDES_PLANNED)
    set_notes(s,
        "Fluorescence is the physical process I exploit. Here's how it works in "
        "five steps. A high-energy photon (UV or short-visible) hits a "
        "molecule and gets absorbed. That kicks an electron up to a higher "
        "energy state. The molecule rapidly loses a little of that energy as "
        "vibration — heat — then the electron drops back down, releasing a NEW "
        "photon. Because of that vibrational loss, the emitted photon has less "
        "energy — longer wavelength — than the photon that excited it. This "
        "wavelength shift is called the Stokes shift. The signature of each "
        "fluorophore is the SHAPE of its emission curve. The right panel lists "
        "some biological fluorophores you can detect this way — NADH and "
        "flavins tell you about cellular metabolism, collagen and elastin "
        "about extracellular matrix, lipofuscin about aging. These are "
        "label-free contrast agents — the body provides them.")

    # ---------- 7. SINGLE-EXCITATION FLUORESCENCE HSI ----------
    s = new_slide()
    add_title_bar(s, "Single-excitation fluorescence HSI")
    # Acquisition chain at top
    add_chip(s, Inches(0.5), Inches(1.2), Inches(2.4), Inches(1.2),
             "UV / blue\nlight source\n(fixed λ_ex)",
             fill=PALE, text_color=ACCENT_BLUE, font_size=12)
    add_text_box(s, Inches(2.95), Inches(1.55), Inches(0.5), Inches(0.5),
                 "→", font_size=28, color=NAVY, align=PP_ALIGN.CENTER)
    add_chip(s, Inches(3.5), Inches(1.2), Inches(2.4), Inches(1.2),
             "Biological\nsample",
             fill=PALE, text_color=ACCENT_GREEN, font_size=12)
    add_text_box(s, Inches(5.95), Inches(1.55), Inches(0.5), Inches(0.5),
                 "→", font_size=28, color=NAVY, align=PP_ALIGN.CENTER)
    add_chip(s, Inches(6.5), Inches(1.2), Inches(2.4), Inches(1.2),
             "HSI camera\nemission @ many λ_em",
             fill=PALE, text_color=ACCENT_RED, font_size=12)
    add_text_box(s, Inches(8.95), Inches(1.55), Inches(0.5), Inches(0.5),
                 "→", font_size=28, color=NAVY, align=PP_ALIGN.CENTER)
    add_chip(s, Inches(9.5), Inches(1.2), Inches(2.4), Inches(1.2),
             "3D cube\n(x, y, λ_em)",
             fill=FAINT, text_color=NAVY, font_size=12)
    # Bottom-left: bullet payoff list
    add_bullet_list(s, Inches(0.5), Inches(2.8), Inches(7.5), Inches(3.8), [
        "Result: at every pixel, a 1D emission spectrum — a chemical fingerprint",
        "Each tissue type shows a characteristic emission curve at the chosen λ_ex",
        "Already enables: tissue typing, NADH-redox mapping, autofluorescence contrast",
        "Compared to RGB:",
        "  - 24+ channels of molecular info per pixel vs. 3 RGB channels",
        "  - Label-free — no exogenous dye injection",
    ], font_size=12)
    # Bottom-right: HSI cube illustration
    add_image(s, IMG3P / "hsi_cube_alt.jpg", Inches(8.4), Inches(2.8),
              width=Inches(4.6))
    add_caption(s, Inches(8.4), Inches(6.55), Inches(4.6),
                "A spectral cube — 1 spectrum / pixel")
    add_image_credit(s, Inches(0.4), Inches(7.05), Inches(12.5),
                     "Hyperspectral cube illustration: Wikimedia Commons CC BY-SA")
    add_footer(s, 7, SLIDES_PLANNED)
    set_notes(s,
        "Single-excitation fluorescence HSI works like this. You shine a "
        "fixed-wavelength UV or blue light source on a biological sample. "
        "Some molecules absorb the light and re-emit fluorescence at longer "
        "wavelengths. A hyperspectral camera captures that emission at many "
        "wavelengths simultaneously. The output is a 3D data cube: two "
        "spatial dimensions and one emission-wavelength dimension. At every "
        "pixel you now have a full emission spectrum — a molecular "
        "fingerprint. Compared to RGB, you've gone from 3 colour channels to "
        "20 or more channels of molecular information per pixel, and crucially "
        "this is label-free — you don't need to inject exogenous dyes. So "
        "single-excitation HSI already enables tissue typing, NADH-redox "
        "mapping, basic autofluorescence contrast. The question is: is that "
        "enough?")

    # ---------- 8. WHY SINGLE-EXCITATION IS NOT ENOUGH ----------
    s = new_slide()
    add_title_bar(s, "Why a single excitation is not enough")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(7.0), Inches(5.5), [
        "Many biological fluorophores overlap in emission",
        "  - Collagen emits near 400 nm, elastin near 410 nm",
        "  - NADH near 460 nm, lipofuscin tails through the same region",
        "From the emission spectrum ALONE, you can't always tell who's emitting",
        "Different fluorophores have different absorption profiles — they respond strongly to DIFFERENT excitations",
        "A single-excitation acquisition picks one slice of this 2D landscape and discards the rest",
        "Result: confused tissue classification when overlapping fluorophores are present",
    ], font_size=15)
    # Right side: the conceptual diagram - two overlapping curves
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(7.8), Inches(1.4), Inches(5.0), Inches(5.0))
    box.fill.solid(); box.fill.fore_color.rgb = PALE
    box.line.color.rgb = NAVY; box.line.width = Pt(1.2)
    add_text_box(s, Inches(7.9), Inches(1.55), Inches(4.8), Inches(0.5),
                 "The disambiguation problem", font_size=15, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(7.9), Inches(2.1), Inches(4.8), Inches(4.2),
                 "Fluorophore A: peaks at ~450 nm\nFluorophore B: peaks at ~460 nm\n\n"
                 "→ Emission spectra OVERLAP heavily.\n\n"
                 "But:\nA absorbs strongly at 340 nm\nB absorbs strongly at 280 nm\n\n"
                 "→ Different EXCITATIONS produce different emission ratios.\n"
                 "    Vary the excitation to disambiguate.",
                 font_size=13, color=DARK_GREY)
    add_footer(s, 8, SLIDES_PLANNED)
    set_notes(s,
        "The problem is that biology stacks fluorophores on top of each other. "
        "Collagen emits near 400 nanometers; elastin near 410; NADH near 460; "
        "lipofuscin tails through the same region. From the emission spectrum "
        "alone — one slice through this landscape — you often can't tell which "
        "fluorophore is responsible for a signal. But different fluorophores "
        "absorb at different wavelengths — they have different excitation "
        "profiles. Even if their emissions overlap, you can disambiguate them "
        "by exciting at different wavelengths and watching how the emission "
        "ratio changes. A single-excitation acquisition discards this "
        "information entirely. The fix — and the focus of this thesis — is to "
        "vary the excitation as well.")

    # ---------- 9. ME-HSI: THE 4D EXTENSION ----------
    s = new_slide()
    add_title_bar(s, "Multi-Excitation HSI: the 4D extension")
    add_bullet_list(s, Inches(0.4), Inches(1.2), Inches(6.5), Inches(5.5), [
        "Acquire the same scene at MANY excitation wavelengths sequentially",
        "Each pixel now has TWO spectral dimensions:",
        "  - emission λ_em (as before)",
        "  - excitation λ_ex (new)",
        "A pixel's data is no longer a curve — it's a 2D excitation-emission matrix (EEM)",
        "Each fluorophore has a characteristic EEM \"signature\"",
        "Overlapping emission curves now separate when you also vary excitation",
    ], font_size=13)
    # Middle: HSI cube illustration (from RGB → multi → hyper)
    add_image(s, IMG3P / "hsi_cube.png", Inches(7.0), Inches(1.2),
              width=Inches(3.3))
    add_caption(s, Inches(7.0), Inches(4.55), Inches(3.3),
                "Single → multi → hyperspectral\n(ME-HSI adds another axis on top)")
    # Right: the 4D structure callout
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(10.45), Inches(1.2), Inches(2.7), Inches(5.3))
    box.fill.solid(); box.fill.fore_color.rgb = PALE
    box.line.color.rgb = ACCENT_RED; box.line.width = Pt(2)
    add_text_box(s, Inches(10.55), Inches(1.35), Inches(2.5), Inches(0.5),
                 "Data shape:", font_size=13, bold=True, color=NAVY)
    add_text_box(s, Inches(10.55), Inches(1.85), Inches(2.5), Inches(0.8),
                 "I(x, y, λ_em, λ_ex)",
                 font_size=15, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(10.55), Inches(2.85), Inches(2.5), Inches(3.5),
                 "Example:\n• 8 excitations\n• 24 em bands each\n• 192 layers / pixel\n• 1040×925 →\n  ~184 M points / frame",
                 font_size=11, color=DARK_GREY)
    add_image_credit(s, Inches(0.4), Inches(7.05), Inches(12.5),
                     "Spectral cube comparison: Wikimedia Commons CC BY-SA")
    add_footer(s, 9, SLIDES_PLANNED)
    set_notes(s,
        "The fix is multi-excitation HSI — ME-HSI. Instead of one fixed "
        "excitation, you sequentially acquire the same scene under MANY "
        "excitation wavelengths. Now each pixel has two spectral dimensions: "
        "emission as before, plus excitation. A single pixel's data is no "
        "longer a 1D curve — it's a 2D excitation-emission matrix, or EEM. "
        "Each fluorophore has a characteristic EEM signature; this is the "
        "gold-standard fluorescence fingerprint used in analytical chemistry. "
        "When two fluorophores overlap in emission, their EEM signatures still "
        "differ. The downside is data volume: for an 8-excitation by "
        "24-emission acquisition, you get 192 spectral layers per pixel. For "
        "a megapixel image, that's nearly 200 million data points per frame.")

    # ---------- 10. ME-HSI CHALLENGES ----------
    s = new_slide()
    add_title_bar(s, "ME-HSI's new challenges")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.5), [
        "Data explosion: 192+ spectral layers per pixel ⇒ acquisition time, storage, computation all scale up",
        "Acquisition is sequential — slow; not ideal for moving samples or live tissue",
        "Physical artifacts:",
        "  - Rayleigh scattering at λ_em ≈ λ_ex creates invalid bands (1st order)",
        "  - 2nd-order harmonics at λ_em ≈ 2·λ_ex add more invalid bands",
        "  - Variable number of valid emission bands per excitation",
        "Spectral redundancy: most (λ_ex, λ_em) cells carry overlapping information",
        "Cross-excitation correlations: the JOINT structure between excitations is precisely what 3D HSI methods cannot see",
        "Detector noise dominates at low-signal regions and long excitations",
    ], font_size=14)
    add_footer(s, 10, SLIDES_PLANNED)
    set_notes(s,
        "ME-HSI is powerful but brings its own challenges. Data explosion is "
        "the first — 192-plus spectral layers per pixel translates to "
        "acquisition delays, storage, and computational cost. Acquisition is "
        "sequential — slow, not ideal for moving samples. Then there are "
        "physical artifacts: Rayleigh scattering at the excitation wavelength "
        "and its second-order harmonic create invalid bands that have to be "
        "masked out, leaving variable numbers of valid emission bands per "
        "excitation. Most of the remaining (excitation, emission) cells carry "
        "overlapping or redundant information. And critically, the JOINT "
        "structure between excitations — which is what makes ME-HSI valuable "
        "in the first place — is exactly what traditional 3D HSI band "
        "selection methods cannot see. So we need analytical tools designed "
        "for the 4D structure. That's the gap this thesis addresses.")

    # ---------- 11. EXISTING BAND SELECTION & WHY THEY FAIL ----------
    s = new_slide()
    add_title_bar(s, "Existing band-selection methods")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.5), [
        "Four standard families in 3D HSI literature:",
        "  - Filter (variance, entropy, correlation): cheap, no task signal",
        "  - Wrapper (sequential forward selection, genetic algorithms): expensive, label-hungry",
        "  - Embedded (LASSO, attention): jointly learn weights and selection — need labels",
        "  - Deep (BS-Net, attention models): learn-to-select via reconstruction",
        "Three common limitations when applied to ME-HSI:",
        "  - Most assume a 3D structure (x, y, λ) — they don't model cross-excitation correlation",
        "  - Most treat bands independently — they miss the joint (λ_ex, λ_em) signal",
        "  - Most need labels — supervised, label-hungry, can't be deployed on unlabeled data",
        "→ Gap: no unsupervised, 4D-aware method that captures cross-excitation correlations",
    ], font_size=14)
    add_footer(s, 11, SLIDES_PLANNED)
    set_notes(s,
        "The band-selection literature is mature for 3D HSI, with four "
        "standard families: filter methods like variance and entropy ranking, "
        "wrapper methods like sequential forward selection or genetic "
        "algorithms, embedded methods like sparse LASSO and attention "
        "mechanisms, and deep methods like BS-Net. All have known issues "
        "though. Most assume a 3D input structure and treat bands "
        "independently — they don't model the joint (excitation, emission) "
        "structure distinctive to ME-HSI. And most either require labels or "
        "carry no task signal at all. There is no published method that's "
        "both unsupervised AND 4D-aware. That's the gap my thesis fills.")

    # ---------- 12. CONTRIBUTIONS ----------
    s = new_slide()
    add_title_bar(s, "Contributions of this thesis")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.5), [
        "Architectural: a 3D convolutional autoencoder for 4D ME-HSI",
        "  - Parallel encoder branches per excitation handle variable emission counts after Rayleigh cutoff",
        "  - Excitation-averaging operator captures cross-excitation correlations in a shared latent representation",
        "Methodological: perturbation-based wavelength attribution",
        "  - Finite-difference Jacobian of reconstruction w.r.t. latent dimensions",
        "  - Produces a causal influence matrix in the (λ_ex, λ_em) space",
        "  - MMR diversity term prevents redundant selections",
        "Empirical: evaluation on Lichens (4 classes, 192 bands) and Collagen Sponges (3 classes, 158 bands)",
        "Validation: robustness vs 10,000 random selections + classifier-family generalization across 10 algorithms",
    ], font_size=14)
    add_footer(s, 12, SLIDES_PLANNED)
    set_notes(s,
        "My contributions are fourfold. First, architectural: a 3D convolutional "
        "autoencoder with parallel encoder branches per excitation. This handles "
        "the variable emission counts after Rayleigh cutoff cleanly, and "
        "excitation-averaging in the shared latent space captures the joint "
        "cross-excitation correlations. Second, methodological: a "
        "perturbation-based attribution that's essentially a finite-difference "
        "Jacobian of the reconstruction with respect to latent dimensions, "
        "combined with Maximum Marginal Relevance for diversity. Third, "
        "empirical evaluation on two distinct datasets — lichens, with four "
        "morphological classes, and collagen sponges, with three crosslinker "
        "concentration classes. Fourth, validation: I tested both the "
        "load-bearing nature of the selection (via 10,000 random comparison "
        "on lichens) and its classifier-family generalisability (via 10 "
        "downstream classifiers on collagen sponges). Together these answer "
        "the four reviewer-grade concerns about unsupervised band selection.")

    # ---------- 13. METHOD OVERVIEW ----------
    s = new_slide()
    add_title_bar(s, "Method overview: three stages")
    stages = [
        ("Stage 1\nRepresentation learning",
         "3D CAE with parallel\nexcitation branches\n→ shared latent z ∈ ℝ²⁰",
         Inches(0.5)),
        ("Stage 2\nAttribution",
         "Latent-space perturbation\nproduces influence matrix\nΔ(λ_ex, λ_em)",
         Inches(4.8)),
        ("Stage 3\nSelection",
         "MMR greedy:\nλ·Rel(w) − (1−λ)·max Sim\n→ top-K bands",
         Inches(9.1)),
    ]
    stage_w = Inches(3.8); stage_h = Inches(2.4)
    for title, body, left in stages:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8),
                                 stage_w, stage_h)
        box.fill.solid(); box.fill.fore_color.rgb = FAINT
        box.line.color.rgb = NAVY; box.line.width = Pt(1.2)
        add_text_box(s, left, Inches(1.9), stage_w, Inches(0.8), title,
                     font_size=16, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, left, Inches(2.75), stage_w, Inches(1.5), body,
                     font_size=13, color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(4.3), Inches(2.95), Inches(0.5), Inches(0.4),
                 "→", font_size=32, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(8.6), Inches(2.95), Inches(0.5), Inches(0.4),
                 "→", font_size=32, color=NAVY, align=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(0.5), Inches(4.7), Inches(12.4), Inches(2.2), [
        "Entirely unsupervised: no labels used in any of the three stages",
        "Output: an ordered list of K (λ_ex, λ_em) pairs — threshold at any K",
        "Same pipeline applies to both datasets without retuning",
    ], font_size=15)
    add_footer(s, 13, SLIDES_PLANNED)
    set_notes(s,
        "The framework has three stages. Stage one is representation learning: "
        "the 3D CAE learns a 20-dimensional shared latent representation of "
        "the 4D cube. Stage two is attribution: I perturb each top latent "
        "dimension and measure the reconstruction sensitivity at each "
        "(excitation, emission) cell — this is a finite-difference "
        "approximation of the Jacobian. Stage three is selection: Maximum "
        "Marginal Relevance balances per-band relevance with diversity, "
        "outputting an ordered list of K bands. Critically, no labels enter any "
        "stage. The same pipeline applies to both datasets without changing a "
        "single hyperparameter.")

    # ---------- 14. STAGE 1: 3D CAE ----------
    s = new_slide()
    add_title_bar(s, "Stage 1: 3D CAE with parallel branches")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(7.0), Inches(5.5), [
        "One Conv3D encoder branch per excitation λ_ex",
        "  - Input: cube X⁽ⁱ⁾ of shape (H, W, N_em(λ_ex))",
        "  - Filter size 5 × 5 × min(5, N_em); k₁ = 20 filters",
        "Adaptive pooling standardises the spectral dim to 1",
        "Excitation averaging → shared latent z ∈ ℝ²⁰",
        "  - h̄ = (1 / N_ex) Σᵢ h⁽ⁱ⁾",
        "  - Captures cross-excitation correlations — the ME-HSI-specific 4D signal",
        "Symmetric Conv3D decoder reconstructs each excitation",
        "Loss: masked MSE; 64×64 patches; Adam at lr=10⁻³; converges in 25-30 epochs",
    ], font_size=14)
    add_text_box(s, Inches(7.9), Inches(1.4), Inches(5.0), Inches(0.5),
                 "Why parallel branches?", font_size=16, bold=True, color=NAVY)
    add_bullet_list(s, Inches(7.9), Inches(1.9), Inches(5.0), Inches(4.5), [
        "Variable N_em per excitation",
        "  - Rayleigh cutoff: each λ_ex has a different valid band count",
        "  - A monolithic kernel would need padding + masking complications",
        "Excitation averaging captures",
        "  - cross-excitation correlations",
        "  - the ME-HSI-specific 4D signal",
        "  - patterns invisible to per-band methods",
    ], font_size=13)
    add_footer(s, 14, SLIDES_PLANNED)
    set_notes(s,
        "Stage one in detail. The architecture is a 3D CAE with one parallel "
        "encoder branch per excitation wavelength. Each branch applies a 3D "
        "convolution with 20 filters; the kernel size is 5×5 spatially and "
        "min-of-5-or-Nem spectrally to handle the variable emission counts "
        "after Rayleigh cutoff. After adaptive pooling each excitation has a "
        "feature map of the same shape, and those are averaged into a shared "
        "20-dimensional latent representation. The key design choice: parallel "
        "branches handle variable emission counts cleanly, and excitation "
        "averaging is what makes this 4D-aware. A monolithic kernel would "
        "need messy padding and masking. The decoder mirrors the encoder, and "
        "we train with masked MSE on 64×64 patches.")

    # ---------- 15. STAGE 2: PERTURBATION ----------
    s = new_slide()
    add_title_bar(s, "Stage 2: Latent-space perturbation")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(2.3), [
        "After training, rank top-k latent dims by variance or PCA loading",
        "For each top-k dim d: perturb the code:  z⁺_d = z + ε·σ_d·e_d  and  z⁻_d = z − ε·σ_d·e_d",
        "Decode both → measure reconstruction sensitivity at every (λ_ex, λ_em)",
    ], font_size=15)
    eq_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.8), Inches(3.7),
                                Inches(9.7), Inches(1.0))
    eq_box.fill.solid(); eq_box.fill.fore_color.rgb = FAINT
    eq_box.line.color.rgb = NAVY
    add_text_box(s, Inches(1.8), Inches(3.75), Inches(9.7), Inches(0.8),
                 "Δ_d(λ_ex, λ_em) = |G(z⁺_d) − G(z⁻_d)| / (2·ε·σ_d)",
                 font_size=18, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, Inches(1.8), Inches(4.7), Inches(9.7), Inches(0.4),
                 "(finite-difference approximation of |∂G/∂z_d|)",
                 font_size=12, italic=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(0.5), Inches(5.3), Inches(12.4), Inches(1.7), [
        "Aggregate across dims weighted by dim score:  Influence(λ_ex, λ_em) = Σ score(d) · mean_ε Δ_d(λ_ex, λ_em)",
        "This is a CAUSAL attribution: it measures what the network actually depends on",
        "Unlike variance ranking, this can find LOW-variance but discriminative bands",
    ], font_size=15)
    add_footer(s, 15, SLIDES_PLANNED)
    set_notes(s,
        "Stage two is the methodological core. The trained autoencoder has a "
        "compressed representation we can probe. I rank the top-k latent "
        "dimensions by variance or PCA loading. For each top dim, I perturb "
        "the code by plus and minus epsilon times sigma_d, decode both, and "
        "measure the reconstruction difference at every (excitation, emission) "
        "cell. This is essentially a finite-difference approximation of the "
        "Jacobian of the decoder with respect to that latent dim. Sum across "
        "dims weighted by their scores and you get the influence matrix. "
        "Crucially this is a CAUSAL attribution — it measures what the network "
        "actually depends on, not just what's variable. Variance ranking would "
        "miss low-variance but discriminative bands; perturbation finds them.")

    # ---------- 16. STAGE 3: MMR ----------
    s = new_slide()
    add_title_bar(s, "Stage 3: Maximum Marginal Relevance selection")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(2.4), [
        "Pure top-K by influence would cluster — adjacent (λ_ex, λ_em) pairs are spectrally redundant",
        "MMR (Carbonell & Goldstein 1998) balances relevance with diversity:",
    ], font_size=15)
    eq_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.0),
                                Inches(10.3), Inches(1.0))
    eq_box.fill.solid(); eq_box.fill.fore_color.rgb = FAINT
    eq_box.line.color.rgb = NAVY
    add_text_box(s, Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.8),
                 "MMR(w) = λ · Rel(w) − (1 − λ) · max_{w' ∈ S} Sim(w, w')",
                 font_size=20, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullet_list(s, Inches(0.5), Inches(4.3), Inches(12.4), Inches(2.5), [
        "Rel(w) = normalized influence from Stage 2",
        "Sim(w, w') = cosine similarity between spectral profiles at the two bands",
        "λ = 0.5 (balanced — pure relevance over-clusters, pure diversity ignores signal)",
        "Greedy selection: pick K bands iteratively, each maximizing MMR over the remainder",
    ], font_size=15)
    add_footer(s, 16, SLIDES_PLANNED)
    set_notes(s,
        "Stage three is selection. If I just took the top-K by influence I'd "
        "get a cluster of highly correlated, adjacent bands. Instead I use "
        "Maximum Marginal Relevance, originally from the information-retrieval "
        "literature. MMR balances per-band relevance against diversity with a "
        "parameter lambda that controls the trade-off. Relevance is the "
        "normalized influence score; diversity is the cosine similarity between "
        "the spectral profiles at the two bands. Lambda = 0.5 gives a balanced "
        "result. Greedy iteration: pick K bands one at a time, each maximizing "
        "MMR over what hasn't been selected yet.")

    # ---------- 17. DATASET 1: LICHENS ----------
    s = new_slide()
    add_title_bar(s, "Dataset 1: Lichens (primary, supervised)")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(6.5), Inches(5.7), [
        "Lichen specimens — fungus + photosynthetic partner symbiotic organism",
        "Autofluorescence reflects both components",
        "Acquisition:",
        "  - 8 excitations: 310, 325, 340, 365, 385, 400, 415, 430 nm",
        "  - Emission: 420 – 720 nm @ 10 nm step",
        "  - LED-based source + Nuance FX camera",
        "After Rayleigh cutoff: 192 valid (λ_ex, λ_em) pairs",
        "Ground truth: 4 morphological classes, ~190k annotated pixels",
        "Spatial: 1040 × 925 pixels",
    ], font_size=15)
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "LichensRGB.jpg",
              Inches(7.3), Inches(1.3), width=Inches(5.5))
    add_caption(s, Inches(7.3), Inches(6.0), Inches(5.5),
                "4 lichen species used in this study")
    add_footer(s, 17, SLIDES_PLANNED)
    set_notes(s,
        "The primary dataset is lichens. Lichens are composite organisms — "
        "fungus plus a photosynthetic partner — and their autofluorescence "
        "captures both components, which makes them a natural fluorescence "
        "test case. We acquired on a custom LED-based source paired with a "
        "Nuance FX hyperspectral camera. Eight excitation wavelengths from 310 "
        "to 430 nanometers, emission from 420 to 720 nanometers at 10 nm "
        "step. After Rayleigh cutoff we have 192 valid (excitation, emission) "
        "pairs. Four morphological classes; about 190,000 annotated pixels in "
        "a 1040 by 925 image. This is our primary benchmark.")

    # ---------- 18. LICHENS GROUND TRUTH ----------
    s = new_slide()
    add_title_bar(s, "Lichens: ground truth and ROI training set")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "LichensLabels and ROI.png",
              Inches(0.5), Inches(1.2), width=Inches(8.5))
    add_bullet_list(s, Inches(9.3), Inches(1.3), Inches(3.6), Inches(5.5), [
        "16 specimens in a 4 × 4 grid",
        "Columns = morphological class",
        "4 small ROIs (top row) = training pixels",
        "Rest of each specimen = test set",
        "  - 11.5k train",
        "  - ~180k test",
        "ROIs used ONLY for the downstream classifier — selection is unsupervised",
    ], font_size=14)
    add_footer(s, 18, SLIDES_PLANNED)
    set_notes(s,
        "Here's the ground truth layout. 16 specimens arranged in a 4 by 4 "
        "grid where columns correspond to the 4 morphological classes. The "
        "small rectangular ROIs in the top row are the supervised training "
        "pixels — about 11,500 of them. The rest of each specimen is the test "
        "set, roughly 180,000 pixels. The ROIs are used only for the downstream "
        "KNN classifier that validates the selection — the selection pipeline "
        "itself sees no labels at any stage.")

    # ---------- 19. DATASET 2: COLLAGEN SPONGES ----------
    s = new_slide()
    add_title_bar(s, "Dataset 2: Collagen Sponges (secondary, supervised)")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(6.5), Inches(5.7), [
        "Collagen sponge samples prepared at 3 crosslinker concentrations",
        "  - Crosslinker concentration modulates fibril organization",
        "  - Produces measurable shifts in emission peak & intensity",
        "Tests cross-domain transfer (lichen → collagen chemistry)",
        "Acquisition:",
        "  - 6 excitations: 310, 325, 340, 365, 385, 400 nm (fewer than Lichens)",
        "  - Same camera + same emission grid",
        "After Rayleigh cutoff: 158 valid (λ_ex, λ_em) pairs",
        "Ground truth: 3 concentration classes, ~40k annotated pixels",
    ], font_size=14)
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "collagen_sponges" / "CollagenLabels_and_ROI.png",
              Inches(7.3), Inches(1.3), width=Inches(5.5))
    add_caption(s, Inches(7.3), Inches(6.0), Inches(5.5),
                "3 × 3 grid: 3 crosslinker classes × 3 replicates")
    add_footer(s, 19, SLIDES_PLANNED)
    set_notes(s,
        "The second dataset is Collagen Sponges. These are collagen sponge "
        "samples at three different crosslinker concentrations — the "
        "crosslinker modulates the fibril organization, producing measurable "
        "shifts in emission peak position and intensity. The key question this "
        "dataset answers is whether the framework, tuned on lichen biology, "
        "transfers to a chemically distinct fluorophore family without "
        "modification. Six excitation channels were used instead of eight — "
        "the longer-UV excitations added negligible signal in preliminary "
        "tests. After Rayleigh cutoff we have 158 valid pairs — a different "
        "grid from Lichens, so the selection task is genuinely new. The "
        "image on the right shows the 3 by 3 specimen layout, with the small "
        "interior ROIs in the top row marking the supervised training pixels.")

    # ---------- 20. LICHENS RESULTS: ACCURACY ENVELOPE ----------
    s = new_slide()
    add_title_bar(s, "Lichens: accuracy vs. number of selected bands")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "accuracy_envelope.png",
              Inches(0.6), Inches(1.2), width=Inches(8.0))
    add_bullet_list(s, Inches(8.9), Inches(1.3), Inches(4.2), Inches(5.5), [
        "Headline results:",
        "  - 192 bands (baseline): 88.2%",
        "  - 80 bands: 95.2% (+7.0 pp, 58% reduction)",
        "  - 13 bands: 90.2% (+2.0 pp)",
        "  - 9 bands: 89.4% (+1.2 pp, 95% reduction)",
        "Above-baseline accuracy from K = 9 up to K ≈ 100",
        "Curve PEAKS at intermediate K — too many bands actively hurts",
    ], font_size=13)
    add_footer(s, 20, SLIDES_PLANNED)
    set_notes(s,
        "Here's the headline Lichens result — the accuracy envelope from the "
        "full 3,072-config sweep. Each shaded band shows the range of "
        "accuracies different attribution configurations achieve at that K, "
        "the solid black line is the mean, and the dashed green line is the "
        "best configuration at each K. The dashed red line is the 192-band "
        "baseline at 88.2 percent. The selected bands actually exceed the "
        "baseline across a wide range. At 80 bands we hit 95.2 percent — 7 "
        "percentage points above baseline at 58 percent data reduction. At 9 "
        "bands we still beat baseline at 89.4 percent with 95 percent data "
        "reduction. The curve actually PEAKS at intermediate K — using too "
        "many bands hurts classification because they add noise.")

    # ---------- 21. LICHENS CLASSIFICATION MAPS ----------
    s = new_slide()
    add_title_bar(s, "Lichens: spatial classification maps")
    map_w = Inches(4.0)
    for (path, left, caption) in [
        ("Baseline.png", Inches(0.3),
         "Baseline (all 192 bands)\n88.2% accuracy"),
        ("80bands-best.png", Inches(4.5),
         "K=80 (58% reduction)\n95.2% accuracy"),
        ("9bands-efficient.png", Inches(8.7),
         "K=9 (95% reduction)\n89.4% accuracy"),
    ]:
        add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / path,
                  left, Inches(1.4), width=map_w)
        add_text_box(s, left, Inches(5.7), map_w, Inches(0.8), caption,
                     font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(6.7), Inches(12.4), Inches(0.4),
                 "Spatial homogeneity improves with selection; specimen boundaries sharpen at K=80",
                 font_size=14, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_footer(s, 21, SLIDES_PLANNED)
    set_notes(s,
        "Same result visualized as classification maps. Each map shows the "
        "pixel-wise predicted class for the 16 specimens. The baseline uses "
        "all 192 bands at 88.2 percent — you can see some boundary noise. The "
        "80-band selection at 95.2 percent has visibly cleaner spatial "
        "structure and sharper specimen boundaries. The 9-band selection at "
        "89.4 percent is still visually competitive with the baseline despite "
        "using 95 percent less data. The visual confirmation matters — the "
        "framework isn't just numerically better, it produces spatially "
        "cleaner outputs.")

    # ---------- 22. LICHENS ROBUSTNESS ----------
    s = new_slide()
    add_title_bar(s, "Lichens: validation against 10,000 random selections")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "robustness_histogram.png",
              Inches(0.6), Inches(1.2), width=Inches(8.0))
    add_bullet_list(s, Inches(8.9), Inches(1.3), Inches(4.2), Inches(5.5), [
        "Compared the 13-band selection vs. 10,000 random 13-band combinations",
        "Random mean: 46.1% (median 44.9%)",
        "Random max: 57.9%",
        "Learned selection: 90.2%",
        "12σ separation from random mean",
        "p ≪ 10⁻²⁰ — NOT chance",
        "Selection IS load-bearing — not a lucky subset",
    ], font_size=13)
    add_footer(s, 22, SLIDES_PLANNED)
    set_notes(s,
        "First validation: are these results just lucky? I sampled 10,000 "
        "random 13-band combinations from the 192 available and applied the "
        "same downstream classifier. The random distribution has a mean of "
        "46.1 percent and a max of 57.9 percent. Our learned selection hits "
        "90.2 percent — 12 standard deviations above the random mean. The "
        "p-value is vanishingly small. The framework is identifying genuinely "
        "informative bands, not lucky combinations.")

    # ---------- 23. LICHENS WAVELENGTH HEATMAP ----------
    s = new_slide()
    add_title_bar(s, "Lichens: which wavelengths matter?")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "wavelength_heatmap.png",
              Inches(0.4), Inches(1.2), width=Inches(8.5))
    add_bullet_list(s, Inches(9.2), Inches(1.3), Inches(3.9), Inches(5.5), [
        "Mean importance across all above-baseline configs",
        "Hot region: 480 – 620 nm emission, 325 – 400 nm excitation",
        "Matches chlorophyll & lichen autofluorescence physics",
        "9-band optimal config selects ONE band per excitation — directly actionable for sensor design",
    ], font_size=13)
    add_footer(s, 23, SLIDES_PLANNED)
    set_notes(s,
        "This heatmap shows mean wavelength importance across all "
        "above-baseline configurations. Bright cells are consistently "
        "top-ranked; dark cells are rarely selected. The hot region spans "
        "roughly 480 to 620 nanometers emission with excitations 325 to 400 "
        "doing most of the work. This matches what we know about lichen "
        "autofluorescence physics — chlorophyll and related fluorophores emit "
        "in this range. And importantly, the 9-band optimal configuration "
        "selects exactly ONE band per excitation source. That's directly "
        "actionable for sensor design: you could build a system with 9 "
        "filters, one per LED, and recover this performance.")

    # ---------- 24. COLLAGEN ACCURACY ENVELOPE ----------
    s = new_slide()
    add_title_bar(s, "Collagen Sponges: accuracy vs. selected bands")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "collagen_sponges" / "accuracy_envelope.png",
              Inches(0.6), Inches(1.2), width=Inches(8.0))
    add_bullet_list(s, Inches(8.9), Inches(1.3), Inches(4.2), Inches(5.5), [
        "158 bands (baseline): 79.8%",
        "5 bands: 80.9% (+1.1 pp)",
        "10 bands: 81.7% (+1.9 pp)",
        "15 bands: 82.4% (+2.6 pp)",
        "20 bands: 85.0% (+5.2 pp)",
        "30 bands: 85.6% (+5.8 pp) — peak",
        "Every K from 5 to 130 exceeds baseline",
        "Same pipeline, no retuning",
    ], font_size=13)
    add_footer(s, 24, SLIDES_PLANNED)
    set_notes(s,
        "Now Collagen Sponges. Same pipeline, no retuning. The 158-band "
        "baseline is 79.8 percent. With just 5 bands we already exceed it at "
        "80.9 percent. The peak is at 30 bands with 85.6 percent — 5.8 "
        "percentage points above baseline at 81 percent data reduction. Every "
        "K from 5 to 130 exceeds the baseline. This is important because it "
        "shows the framework transfers across sample chemistries — lichen and "
        "collagen autofluorescence have different fluorophore populations, "
        "but the same selection method works on both.")

    # ---------- 25. COLLAGEN WAVELENGTH HEATMAP ----------
    s = new_slide()
    add_title_bar(s, "Collagen Sponges: wavelength importance map")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "collagen_sponges" / "wavelength_heatmap.png",
              Inches(0.4), Inches(1.5), width=Inches(8.5))
    add_bullet_list(s, Inches(9.2), Inches(1.5), Inches(3.9), Inches(5.0), [
        "6 excitations × emission bands",
        "Hot region:",
        "  - 600 – 690 nm emission",
        "  - 365 – 400 nm excitation",
        "Matches collagen autofluorescence physics:",
        "  - Crosslinker-modulated AGE emission",
        "  - Native collagen near short-wave edge",
        "Grey: Rayleigh-invalid cells",
    ], font_size=13)
    add_footer(s, 25, SLIDES_PLANNED)
    set_notes(s,
        "Same wavelength importance heatmap structure as Lichens but for "
        "Collagen Sponges. Bright = consistently top-ranked across "
        "configurations. The hot region concentrates at 600 to 690 nanometers "
        "emission under 365 to 400 nanometer excitation. There's a secondary "
        "band at the short-wavelength edge across all excitations. This "
        "matches the chemistry: the long-wavelength region captures advanced "
        "glycation end-product cross-link emission, which scales with "
        "crosslinker concentration; the short-wave shoulder is native collagen "
        "fluorescence. Grey cells are Rayleigh-invalid positions.")

    # ---------- 26. COLLAGEN MULTI-CLASSIFIER ----------
    s = new_slide()
    add_title_bar(s, "Collagen Sponges: classifier-family generalization")
    add_image(s, ROOT / "MasterThesis_Narek_Meloyan" / "figures" / "collagen_sponges" / "classifier_curves.png",
              Inches(0.5), Inches(1.2), width=Inches(8.0))
    add_bullet_list(s, Inches(8.7), Inches(1.3), Inches(4.5), Inches(5.5), [
        "Same selection, 10 different classifiers:",
        "  - KNN (5, 11) — uniform + distance",
        "  - LDA",
        "  - SVM (linear, RBF)",
        "  - MLP",
        "  - Random Forest (100, 300)",
        "  - Gradient Boosting",
        "9 of 10 classifiers benefit",
        "LDA peaks: 92.5% on K=50 (vs. 92.8% on all 158)",
        "Selected bands carry PORTABLE information — not KNN-tuned",
    ], font_size=12)
    add_footer(s, 26, SLIDES_PLANNED)
    set_notes(s,
        "Second validation: classifier-family generalization. To test whether "
        "the gains depend on the specific KNN classifier used for scoring, I "
        "evaluated ten different classifiers on the same selection — two KNN "
        "variants, LDA, two SVM kernels, MLP, two Random Forest sizes, and "
        "Gradient Boosting. Nine of ten benefit from the selection. The most "
        "striking result: Linear Discriminant Analysis at K=50 hits 92.5 "
        "percent accuracy, which is within 0.3 percentage points of the 92.8 "
        "percent it achieves on all 158 bands. The linearly separable "
        "structure of the data is essentially preserved by the selection. The "
        "bands are intrinsically informative, not tuned to one classifier.")

    # ---------- 27. DISCUSSION ----------
    s = new_slide()
    add_title_bar(s, "Discussion: what does this mean?")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.7), [
        "Across two ME-HSI datasets (lichens, collagen sponges):",
        "  - Selected wavelengths consistently MEET or EXCEED the full-spectrum baseline",
        "  - At 80 – 95% data reduction — i.e. fewer channels, faster inference",
        "Three reasons it works:",
        "  - Causal attribution: perturbation measures what the network DEPENDS on, not just what's variable",
        "  - Joint 4D representation: cross-excitation correlations captured by parallel branches + averaging",
        "  - Diversity-aware MMR: prevents redundant clustered selections",
        "Cross-domain transfer:",
        "  - Same pipeline, no retuning, works on two chemically distinct samples",
        "  - Only the dim-scoring strategy (PCA-k=3 for Lichens, variance-k=1 for Collagen) needs per-dataset choice",
    ], font_size=14)
    add_footer(s, 27, SLIDES_PLANNED)
    set_notes(s,
        "So what does this mean. Across both datasets, the selected "
        "wavelengths meet or exceed the full-spectrum baseline at 80 to 95 "
        "percent data reduction. Three reasons this works. First, perturbation "
        "gives a causal attribution — it measures what the network actually "
        "depends on, unlike variance which just measures variability. Second, "
        "the parallel-branch architecture with averaging is what makes this "
        "4D-aware — it captures cross-excitation correlations that per-band "
        "methods miss. Third, MMR's diversity term prevents the kind of "
        "clustered selections that pure relevance-ranking produces. The same "
        "pipeline works on both datasets without retuning. The only meaningful "
        "per-dataset knob is the dimension-scoring strategy — Lichens prefers "
        "PCA with k=3, Collagen Sponges prefers variance with k=1.")

    # ---------- 28. LIMITATIONS ----------
    s = new_slide()
    add_title_bar(s, "Limitations & honest scope")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.7), [
        "Computational cost:",
        "  - Full pipeline ~2 hr on RTX 4070 for Lichens (~30 min train + ~90 min sweep)",
        "  - 75% of compute is the perturbation phase — embarrassingly parallel, ~3× speedup possible",
        "  - One-time per sample type — amortised over downstream inference (K/192 ≈ 21× faster at K=9)",
        "Black-box attribution:",
        "  - Perturbation tells WHICH bands matter, not WHY (in chemical terms)",
        "  - Linking selections to specific fluorophores still requires domain expertise",
        "Threshold selection:",
        "  - User must specify K — automated elbow detection is future work",
        "Dataset breadth:",
        "  - Two ME-HSI fluorescence datasets; extending to remote-sensing HSI is future work",
    ], font_size=14)
    add_footer(s, 28, SLIDES_PLANNED)
    set_notes(s,
        "Honest limitations. First, computational cost: training plus the "
        "perturbation sweep takes about 2 hours on our RTX 4070 for Lichens. "
        "Seventy-five percent of that is the perturbation phase — which is "
        "embarrassingly parallel, and an optimised implementation could shave "
        "another 3x. The cost is one-time per sample type, and it's amortised "
        "over all downstream inferences — those are about 21 times faster at "
        "K=9. Second, the attribution is black-box in the chemical sense — "
        "perturbation tells which bands matter, not always why; mapping our "
        "selections to specific fluorophores still requires domain expertise. "
        "Third, the user has to specify K — we don't auto-detect the elbow. "
        "Fourth, dataset breadth: I evaluated on two ME-HSI fluorescence "
        "datasets here; extending to remote-sensing HSI is future work.")

    # ---------- 29. CONCLUSION ----------
    s = new_slide()
    add_title_bar(s, "Conclusion & future directions")
    add_bullet_list(s, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.7), [
        "Summary:",
        "  - First deep-learning wavelength selection method designed for 4D ME-HSI",
        "  - Three stages: parallel-branch 3D CAE + perturbation attribution + MMR diversity",
        "  - Lichens: 95.2% @ K=80 (108% of baseline) | 89.4% @ K=9 (95% reduction)",
        "  - Collagen Sponges: 85.6% @ K=30 (107% of baseline)",
        "  - Validated via robustness vs 10,000 random AND classifier-family generalization (10 classifiers)",
        "Future work:",
        "  - Larger ME-HSI corpus: more sample types, more excitation grids",
        "  - Automated K selection via reconstruction-knee or stability criteria",
        "  - Interpretable attribution: link selected bands to fluorophore databases",
        "  - Uncertainty quantification via variational autoencoders",
    ], font_size=14)
    add_footer(s, 29, SLIDES_PLANNED)
    set_notes(s,
        "To conclude. The contribution: the first deep-learning wavelength "
        "selection method designed specifically for 4D ME-HSI. Three stages: a "
        "parallel-branch convolutional autoencoder, perturbation-based "
        "attribution, and MMR diversity. On Lichens, 95 percent accuracy with "
        "80 of 192 bands, or 89 percent with just 9 of 192. On Collagen "
        "Sponges, 86 percent with 30 of 158 bands. Validated via two "
        "complementary procedures — random-baseline robustness on Lichens and "
        "ten-classifier generalization on Collagen Sponges. Future directions "
        "include larger datasets, automated K selection, fluorophore-level "
        "interpretation, and uncertainty quantification via variational "
        "autoencoders.")

    # ---------- 30. IMAGE CREDITS ----------
    s = new_slide()
    add_title_bar(s, "Image credits & acknowledgments")
    add_text_box(s, Inches(0.5), Inches(1.1), Inches(12.4), Inches(0.5),
                 "Third-party images used in this presentation",
                 font_size=16, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(4.6), [
        "All third-party images sourced from Wikimedia Commons under Creative Commons licenses",
        "Slide 3 — vision in nature:",
        "  - Honeybee on flower; Male monarch butterfly; Mantis shrimp (Gilli Banta reef); EM spectrum band",
        "Slide 4 — HSI applications:  Landsat-5 false-colour, Ethiopia (NASA / USGS, public domain)",
        "Slide 5 — biomedicine:  GFP super-resolution micrograph, Christoph Cremer",
        "Slide 6 — fluorescence:  Jablonski diagram; Fluorescent minerals under UV, Hannes Grobe",
        "Slides 1, 7, 9 — HSI cubes:  Mono / Multi / Hyperspectral spectral signatures illustration",
        "Full attribution table: MasterThesis_Narek_Meloyan/images_third_party/CREDITS.md",
    ], font_size=12)
    add_text_box(s, Inches(0.5), Inches(6.0), Inches(12.4), Inches(0.5),
                 "Project figures (datasets, results, classification maps) are original "
                 "work generated from the spectral-select pipeline.",
                 font_size=12, italic=True, color=DARK_GREY)
    add_footer(s, 30, SLIDES_PLANNED)
    set_notes(s,
        "Quick acknowledgment slide. All the conceptual / vision-biology / "
        "fluorescence-physics images used in the introduction are from "
        "Wikimedia Commons under Creative Commons licenses; the satellite "
        "image is NASA / USGS public domain. Project figures — the datasets, "
        "results, classification maps, accuracy envelopes, wavelength heatmaps "
        "— are all original work generated by the spectral-select pipeline. "
        "Full attribution table lives in the repository.")

    # ---------- 31. THANK YOU ----------
    s = new_slide()
    add_text_box(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.0),
                 "Thank you",
                 font_size=64, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.6),
                 "Questions?",
                 font_size=32, italic=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
                 "Code: github.com/Nard248/spectral-select",
                 font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.4),
                 "narek_meloyan@edu.aua.am",
                 font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    set_notes(s,
        "Thank you for your attention. I'd be happy to take questions. "
        "Code is available on GitHub if anyone wants to look at the "
        "implementation; my email is on the slide for follow-up.")

    prs.save(OUT)
    print(f"Wrote {OUT.relative_to(ROOT)}  ({OUT.stat().st_size / 1024:.0f} KB)")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
